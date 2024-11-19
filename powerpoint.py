import collections, os 
import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PIL import ImageColor
directC=os.path.realpath(os.path.dirname(__file__)) # Current directory
directD=os.path.join(directC,"Data") # Data directory

# EXCLUDE = {"VT01": None,    # Excludes entire patient
#            "VT02": [1,2,3], # Excludes specific sequences of that patient
#            "VT03": ...}
EXCLUDE = {} # Exlude specific atients and/or sequences (if they are already correctly segmented)




## -----------------
colors = {"orange"   : [RGBColor(*ImageColor.getcolor('#EE7733', "RGB")), RGBColor(*ImageColor.getcolor('#EE8866', "RGB"))],
          "grey"     : [RGBColor(*ImageColor.getcolor('#BBBBBB', "RGB")), RGBColor(*ImageColor.getcolor('#FFFFFF', "RGB"))],
          "teal"     : [RGBColor(*ImageColor.getcolor('#009988', "RGB")), RGBColor(*ImageColor.getcolor('#EE7733', "RGB"))],
          "blue"     : [RGBColor(*ImageColor.getcolor('#0077BB', "RGB")), RGBColor(*ImageColor.getcolor('#77AADD', "RGB"))]}
color_loop = [rgb for color, rgb in colors.items()]

def append_pptx(prs, images, color, title=''):
    titleslide = prs.slides.add_slide(prs.slide_layouts[0]) # Specify which patient and sequence
    titleslide.shapes.title.text = title
    titleslide.background.fill.solid()
    titleslide.background.fill.fore_color.rgb = color[0]

    dI = Inches(0.1)
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(image, 0+dI, 0+dI, Inches(10)-2*dI, Inches(7.5)-2*dI) # size of 4:3 slide
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color[1]
    return prs

def generate_pptx(data_dir, exclude, colors, output_dir=''):
    prs = Presentation() # Initialize the powerpoint
    c = 0

    folders = [folder.name for folder in os.scandir(data_dir)]
    for folder in folders:
        if c == len(color_loop): c = 0 # reset coloring

        patient_folder = os.path.join(data_dir,folder)
        if os.path.split(patient_folder)[-1] in exclude.keys():
            excl_seq = exclude[os.path.split(patient_folder)[-1]] # Check which sequences are excluded    
        else:
            excl_seq = []
        sequences = [sequence.name for sequence in os.scandir(patient_folder)]

        for sequence in sequences:
            if excl_seq == None or int(sequence[3:]) in excl_seq:
                break 
            path = os.path.join(patient_folder,sequence,"images")
            if os.path.exists(path):
                images = [os.path.join(path,img) for img in os.scandir(path) if os.path.splitext(img)[1]==".png"]
                append_pptx(prs, images, colors[c], title=f"Patient {folder} seq. {sequence[3:]}")    
        c += 1

    
    prs.save(os.path.join(output_dir,'segmentation_overview.pptx'))
    return

if __name__ == '__main__':
   generate_pptx(directD, EXCLUDE, color_loop)

